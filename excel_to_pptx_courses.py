import openpyxl
from pptx import Presentation
from pptx.util import Pt, Inches
import os

# Function to set font size
def set_font_size(shape, size):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)

# Function to add an image to a slide
def add_image_to_slide(slide, image_path):
    left = Inches(9.0)  # Position from the left of the slide
    top = Inches(5.0)  # Position from the top of the slide
    height = Inches(2.0)  # Height of the image
    slide.shapes.add_picture(image_path, left, top, height=height)

# Function to process each sheet
def process_sheet(sheet, prs):
    # Read the column headers
    print("**********************  " + sheet.title)
    headers = [cell.value for cell in sheet[1]]
    column_indices = {
        'Project': headers.index('Module'),
        'Project description': headers.index('Topics'),
        'Module': headers.index('Detail'),
        'Module Description': headers.index('Examples & Use Cases'),
        'Image': headers.index('Image')  # Index for image column
    }

    # Loop through the data
    for row in sheet.iter_rows(min_row=2, values_only=True):  # Iterating through the rows
        project = row[column_indices['Project']]
        project_description = row[column_indices['Project description']]
        module = row[column_indices['Module']]
        module_description = row[column_indices['Module Description']]
        image_link = row[column_indices['Image']]

        # Construct the full path to the image
        image_path = os.path.join('images', image_link) if image_link else None
        
        print(project, project_description, module, module_description, f"Image path: {image_path}")

        if project == 'Done':
            break

        # Create the initial slide with Project Description
        if project:
            # Create a main heading slide for the project
            slide_layout = prs.slide_layouts[0]  # Layout 0 is Title Slide
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            
            title.text = project
            # No content placeholder used

            # Set font size for the title
            set_font_size(title, 60)

        # Create a subsequent slide for each module
        if project_description not in ['Practice Lab', 'Quiz']:
            if module:
                module_slide_layout = prs.slide_layouts[1]  # Assuming 1 is the Title and Content layout
                module_slide = prs.slides.add_slide(module_slide_layout)
                title = module_slide.shapes.title
                content = module_slide.placeholders[1]
                title.text = f"{project_description}"
                content.text = f"Detail: {module}\n\nExamples and Use Cases: {module_description}"

                # Set font size for title and content
                set_font_size(title, 30)  # Font size for module title
                set_font_size(content, 24)  # Font size for module content

                # Add image to the module slide if the path exists
                if image_path and os.path.exists(image_path):
                    add_image_to_slide(module_slide, image_path)
                else:
                    print(f"Image not found for {image_link}")

# Function to remove the first slide from the presentation
def remove_first_slide(prs):
    if len(prs.slides) > 0:
        xml_slides = prs.slides._sldIdLst  # Access the XML slide list
        slide_id = xml_slides[0]  # Get the ID of the first slide
        xml_slides.remove(slide_id)  # Remove the first slide ID from the list

# Open the Excel file
wb = openpyxl.load_workbook('D:\\Projects\\utilities\\Data Analytics LJ - Final.xlsx')

# Filter out sheets that should not be processed
sheets_to_process = [sheet_name for sheet_name in wb.sheetnames 
                     if "Quiz" not in sheet_name and "topic" not in sheet_name]

# Process each sheet in the filtered list
for sheet_name in sheets_to_process:
    sheet = wb[sheet_name]

    # Load the sample PowerPoint to use its theme
    prs = Presentation('DATAASTRAA_SAMPLE_PPTX.pptx')

    # Process the sheet
    process_sheet(sheet, prs)

    # Remove the first slide from the presentation
    remove_first_slide(prs)

    # Define output path for the current sheet
    output_pptx_path = f'D:\\Projects\\utilities\\Modules\\project_modules_presentation_{sheet_name}.pptx'

    # Ensure the directory exists
    output_dir = os.path.dirname(output_pptx_path)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Save the presentation
    try:
        prs.save(output_pptx_path)
        print(f"Presentation for sheet '{sheet_name}' saved as '{output_pptx_path}'")
    except PermissionError:
        print(f"PermissionError: Unable to save the presentation to '{output_pptx_path}'. Check if the file is open or if you have write permissions.")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")
