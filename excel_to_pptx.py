import openpyxl
from pptx import Presentation
from pptx.util import Inches

# Open the Excel file
wb = openpyxl.load_workbook('Client Use Cases.xlsx')
sheet = wb.active

# Load the sample PowerPoint to use its theme
prs = Presentation('DATAASTRAA_SAMPLE_PPTX.pptx')

current_project = ''

# Loop through the data
for row in sheet.iter_rows(min_row=2, values_only=True):  # Assuming data starts from the second row
    project, project_description, module, module_description = row
    print(project, project_description, module, module_description)

    if project == 'Done':
        break

    if project:
        current_project = project
        # Create a main heading slide for the project
        slide_layout = prs.slide_layouts[1]  # Assuming 1 is the Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = project
        content.text = project_description

    # Create a subsequent slide for each module
    module_slide_layout = prs.slide_layouts[1]  # Assuming 1 is the Title and Content layout
    module_slide = prs.slides.add_slide(module_slide_layout)
    title = module_slide.shapes.title
    content = module_slide.placeholders[1]
    title.text = module
    content.text = module_description

print('for loop done')
# Save the presentation
output_pptx_path = 'project_modules_presentation.pptx'
prs.save(output_pptx_path)

print(f"Presentation saved as '{output_pptx_path}'")



# import os
# import openpyxl
# from pptx import Presentation
# from pptx.util import Inches, Pt

# # Open the Excel file
# wb = openpyxl.load_workbook('Client Use Cases.xlsx')
# sheet = wb.active

# # Load the sample PowerPoint to use its theme
# prs = Presentation('DATAASTRAA_SAMPLE_PPTX.pptx')

# # Extract layout, position, and size from the sample slides
# project_sample_slide = prs.slides[0]
# module_sample_slide = prs.slides[1]



# # Get image position and size from the project sample slide
# project_image_placeholder = project_sample_slide.placeholders[0]
# project_image_left = project_image_placeholder.left
# project_image_top = project_image_placeholder.top
# project_image_width = project_image_placeholder.width
# project_image_height = project_image_placeholder.height


# # print(module_sample_slide.placeholders[3])

# # # Get image position and size from the module sample slide
# # module_image_placeholder = module_sample_slide.placeholders[0]
# # module_image_left = module_image_placeholder.left
# # module_image_top = module_image_placeholder.top
# # module_image_width = module_image_placeholder.width
# # module_image_height = module_image_placeholder.height

# # # Path to the folder containing images
# # images_folder = 'path/to/your/images/folder'

# # current_project = ''

# # # Loop through the data
# # for row in sheet.iter_rows(min_row=2, values_only=True):  # Assuming data starts from the second row
# #     project, project_description, module, module_description = row
# #     # print(project, project_description, module, module_description)

# #     if project == 'Done':
# #         break

# #     if project:
# #         current_project = project
# #         # Create a main heading slide for the project
# #         slide_layout = prs.slide_layouts[5]  # Assuming 5 is the Title and Content layout
# #         slide = prs.slides.add_slide(slide_layout)
# #         title = slide.shapes.title
# #         content = slide.placeholders[0]
# #         title.text = project
# #         content.text = project_description

# #         # Add an image to the project slide if it exists
# #         image_path = os.path.join(images_folder, f"{project}.jpg")  # Assuming images are named after the project
# #         if os.path.exists(image_path):
# #             slide.shapes.add_picture(image_path, project_image_left, project_image_top, project_image_width, project_image_height)

# #     # Create a subsequent slide for each module
# #     module_slide_layout = prs.slide_layouts[6]  # Assuming 6 is the Title and Content layout
# #     module_slide = prs.slides.add_slide(module_slide_layout)
# #     title = module_slide.shapes.title
# #     content = module_slide.placeholders[0]
# #     title.text = module
# #     content.text = module_description

# #     # Add an image to the module slide if it exists
# #     image_path = os.path.join(images_folder, f"{module}.jpg")  # Assuming images are named after the module
# #     if os.path.exists(image_path):
# #         module_slide.shapes.add_picture(image_path, module_image_left, module_image_top, module_image_width, module_image_height)

# # print('for loop done')

# # # Save the presentation
# # output_pptx_path = 'project_modules_presentation.pptx'
# # prs.save(output_pptx_path)

# # print(f"Presentation saved as '{output_pptx_path}'")
